"""
Document Loader Service

Handles all document ingestion logic:
- File type detection
- Full text extraction from various document types
- OCR processing when enabled
- Page-level text extraction for multi-page documents

Supported file types:
- PDF: Extract text from all pages, OCR scanned pages if enabled
- TXT: Read full file content (UTF-8 with fallback encodings)
- DOCX: Extract paragraphs, tables, and headers
- XLSX: Extract cell values from all sheets
- PPTX: Extract text from all slides
- Images (PNG, JPG, JPEG): OCR if enabled, otherwise empty text

All text extraction preserves UTF-8 encoding for proper Arabic support.
"""

import os
import mimetypes
from typing import Dict, Any, List, Optional
from dataclasses import dataclass

from pypdf import PdfReader

# Optional dependencies - gracefully handle if not installed
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from app.config import settings


@dataclass
class PageContent:
    """Content extracted from a single page/section."""
    page_number: int
    text: str
    ocr_applied: bool


@dataclass
class ExtractionResult:
    """Result of document text extraction."""
    full_text: str
    page_count: int
    pages: List[PageContent]
    ocr_applied: bool
    metadata: Dict[str, Any]
    error: Optional[str] = None


class OCRNotAvailableError(Exception):
    """Raised when OCR is requested but Tesseract is not installed."""
    pass


class UnsupportedFileTypeError(Exception):
    """Raised when file type is not supported."""
    pass


class DocumentLoader:
    """
    Handles text extraction from various document types.

    Usage:
        loader = DocumentLoader(enable_ocr=True)
        result = loader.extract(file_path)
    """

    # Supported MIME types and their handlers
    SUPPORTED_TYPES = {
        # PDF
        "application/pdf": "pdf",
        # Text
        "text/plain": "txt",
        # Images
        "image/png": "image",
        "image/jpeg": "image",
        "image/jpg": "image",
        # Microsoft Office
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        # Legacy Office (map to new handlers)
        "application/msword": "docx",
        "application/vnd.ms-excel": "xlsx",
        "application/vnd.ms-powerpoint": "pptx",
    }

    # Extension to MIME type mapping for reliable detection
    EXTENSION_MAP = {
        ".pdf": "application/pdf",
        ".txt": "text/plain",
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".doc": "application/msword",
        ".xls": "application/vnd.ms-excel",
        ".ppt": "application/vnd.ms-powerpoint",
    }

    def __init__(self, enable_ocr: bool = False):
        self.enable_ocr = enable_ocr
        self._ocr_service = None
        self._ocr_checked = False
        self._ocr_available = False

    def _get_ocr_service(self):
        """Lazy-load OCR service only when needed."""
        if self._ocr_service is not None:
            return self._ocr_service

        if not self._ocr_checked:
            self._ocr_checked = True
            try:
                from app.services.ocr_service import OCRService
                self._ocr_service = OCRService()
                self._ocr_available = True
            except Exception as e:
                self._ocr_available = False
                if self.enable_ocr:
                    raise OCRNotAvailableError(str(e))

        if self.enable_ocr and not self._ocr_available:
            raise OCRNotAvailableError(
                "OCR is not available. Please install Tesseract:\n"
                "1. Download from: https://github.com/UB-Mannheim/tesseract/wiki\n"
                "2. Install and note the path\n"
                "3. Set TESSERACT_PATH in your .env file"
            )

        return self._ocr_service

    def detect_mime_type(self, file_path: str) -> str:
        """Detect the MIME type of a file."""
        # First try by extension (more reliable for Office files)
        ext = os.path.splitext(file_path)[1].lower()
        if ext in self.EXTENSION_MAP:
            return self.EXTENSION_MAP[ext]

        # Fallback to mimetypes library
        mime_type, _ = mimetypes.guess_type(file_path)
        return mime_type or "application/octet-stream"

    def extract(self, file_path: str) -> ExtractionResult:
        """
        Extract text from a document.

        Args:
            file_path: Path to the file to process

        Returns:
            ExtractionResult with extracted text and metadata
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        mime_type = self.detect_mime_type(file_path)
        handler_type = self.SUPPORTED_TYPES.get(mime_type)

        if not handler_type:
            supported = "PDF, TXT, DOCX, XLSX, PPTX, PNG, JPG, JPEG"
            raise UnsupportedFileTypeError(
                f"Unsupported file type: {mime_type}. Supported types: {supported}"
            )

        # Route to appropriate handler
        handlers = {
            "pdf": self._extract_pdf,
            "txt": self._extract_txt,
            "image": self._extract_image,
            "docx": self._extract_docx,
            "xlsx": self._extract_xlsx,
            "pptx": self._extract_pptx,
        }

        handler = handlers.get(handler_type)
        if handler:
            return handler(file_path)

        raise UnsupportedFileTypeError(f"No handler for type: {handler_type}")

    def _extract_pdf(self, file_path: str) -> ExtractionResult:
        """Extract text from PDF, with OCR fallback for scanned pages."""
        reader = PdfReader(file_path)
        pages: List[PageContent] = []
        text_parts: List[str] = []
        any_ocr_applied = False
        metadata: Dict[str, Any] = {}

        # Extract PDF metadata
        if reader.metadata:
            if reader.metadata.title:
                metadata["title"] = reader.metadata.title
            if reader.metadata.author:
                metadata["author"] = reader.metadata.author
            if reader.metadata.creation_date:
                metadata["creation_date"] = str(reader.metadata.creation_date)

        total_pages = len(reader.pages)

        for page_num in range(total_pages):
            page = reader.pages[page_num]
            page_text = page.extract_text() or ""
            page_text = page_text.strip()
            page_ocr_applied = False

            # If no text and OCR is enabled, try OCR
            if not page_text and self.enable_ocr:
                try:
                    ocr_service = self._get_ocr_service()
                    page_text = ocr_service.ocr_pdf_page(file_path, page_num)
                    page_text = page_text.strip()
                    page_ocr_applied = True
                    any_ocr_applied = True
                except OCRNotAvailableError:
                    raise
                except Exception:
                    page_text = ""

            pages.append(PageContent(
                page_number=page_num + 1,
                text=page_text,
                ocr_applied=page_ocr_applied
            ))

            if page_text:
                text_parts.append(f"[Page {page_num + 1}]\n{page_text}")

        full_text = "\n\n".join(text_parts)

        return ExtractionResult(
            full_text=full_text,
            page_count=total_pages,
            pages=pages,
            ocr_applied=any_ocr_applied,
            metadata=metadata
        )

    def _extract_txt(self, file_path: str) -> ExtractionResult:
        """Extract text from a plain text file."""
        # Try multiple encodings for Arabic support
        encodings = ["utf-8", "utf-8-sig", "utf-16", "cp1256", "iso-8859-6", "latin-1", "cp1252"]

        content = ""
        used_encoding = None
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    content = f.read()
                used_encoding = encoding
                break
            except UnicodeDecodeError:
                continue

        line_count = content.count("\n") + 1 if content else 0

        return ExtractionResult(
            full_text=content,
            page_count=1,
            pages=[PageContent(page_number=1, text=content, ocr_applied=False)],
            ocr_applied=False,
            metadata={
                "line_count": line_count,
                "character_count": len(content),
                "encoding": used_encoding
            }
        )

    def _extract_docx(self, file_path: str) -> ExtractionResult:
        """Extract text from Word document (.docx)."""
        if not DOCX_AVAILABLE:
            raise UnsupportedFileTypeError(
                "DOCX support requires python-docx. Install with: pip install python-docx"
            )

        doc = DocxDocument(file_path)
        text_parts: List[str] = []
        metadata: Dict[str, Any] = {}

        # Extract document properties
        if doc.core_properties:
            if doc.core_properties.title:
                metadata["title"] = doc.core_properties.title
            if doc.core_properties.author:
                metadata["author"] = doc.core_properties.author
            if doc.core_properties.created:
                metadata["created"] = str(doc.core_properties.created)

        # Extract paragraphs
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                text_parts.append(text)

        # Extract tables
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    table_rows.append(" | ".join(cells))
            if table_rows:
                text_parts.append("\n".join(table_rows))

        full_text = "\n\n".join(text_parts)

        metadata["paragraph_count"] = len(doc.paragraphs)
        metadata["table_count"] = len(doc.tables)

        return ExtractionResult(
            full_text=full_text,
            page_count=1,  # DOCX doesn't have page concept until rendered
            pages=[PageContent(page_number=1, text=full_text, ocr_applied=False)],
            ocr_applied=False,
            metadata=metadata
        )

    def _extract_xlsx(self, file_path: str) -> ExtractionResult:
        """Extract text from Excel spreadsheet (.xlsx)."""
        if not XLSX_AVAILABLE:
            raise UnsupportedFileTypeError(
                "XLSX support requires openpyxl. Install with: pip install openpyxl"
            )

        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text_parts: List[str] = []
        pages: List[PageContent] = []
        metadata: Dict[str, Any] = {}

        sheet_names = workbook.sheetnames
        metadata["sheet_count"] = len(sheet_names)
        metadata["sheet_names"] = sheet_names

        for sheet_idx, sheet_name in enumerate(sheet_names):
            sheet = workbook[sheet_name]
            sheet_text_parts = [f"[Sheet: {sheet_name}]"]

            for row in sheet.iter_rows():
                row_values = []
                for cell in row:
                    if cell.value is not None:
                        # Convert to string and strip
                        cell_text = str(cell.value).strip()
                        if cell_text:
                            row_values.append(cell_text)
                if row_values:
                    sheet_text_parts.append(" | ".join(row_values))

            sheet_text = "\n".join(sheet_text_parts)
            text_parts.append(sheet_text)

            pages.append(PageContent(
                page_number=sheet_idx + 1,
                text=sheet_text,
                ocr_applied=False
            ))

        full_text = "\n\n".join(text_parts)

        return ExtractionResult(
            full_text=full_text,
            page_count=len(sheet_names),
            pages=pages,
            ocr_applied=False,
            metadata=metadata
        )

    def _extract_pptx(self, file_path: str) -> ExtractionResult:
        """Extract text from PowerPoint presentation (.pptx)."""
        if not PPTX_AVAILABLE:
            raise UnsupportedFileTypeError(
                "PPTX support requires python-pptx. Install with: pip install python-pptx"
            )

        prs = Presentation(file_path)
        text_parts: List[str] = []
        pages: List[PageContent] = []
        metadata: Dict[str, Any] = {}

        slide_count = len(prs.slides)
        metadata["slide_count"] = slide_count

        for slide_idx, slide in enumerate(prs.slides):
            slide_text_parts = [f"[Slide {slide_idx + 1}]"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        slide_text_parts.append(text)

                # Extract text from tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text_parts.append(" | ".join(row_text))

            slide_text = "\n".join(slide_text_parts)
            text_parts.append(slide_text)

            pages.append(PageContent(
                page_number=slide_idx + 1,
                text=slide_text,
                ocr_applied=False
            ))

        full_text = "\n\n".join(text_parts)

        return ExtractionResult(
            full_text=full_text,
            page_count=slide_count,
            pages=pages,
            ocr_applied=False,
            metadata=metadata
        )

    def _extract_image(self, file_path: str) -> ExtractionResult:
        """Extract text from an image using OCR."""
        metadata: Dict[str, Any] = {}

        # Get image dimensions
        if PIL_AVAILABLE:
            try:
                with Image.open(file_path) as img:
                    metadata["width"] = img.width
                    metadata["height"] = img.height
                    metadata["format"] = img.format
            except Exception:
                pass

        # If OCR is not enabled, return empty text with warning
        if not self.enable_ocr:
            return ExtractionResult(
                full_text="",
                page_count=1,
                pages=[PageContent(page_number=1, text="", ocr_applied=False)],
                ocr_applied=False,
                metadata=metadata,
                error="OCR not enabled - no text extracted from image"
            )

        # Perform OCR
        try:
            ocr_service = self._get_ocr_service()
            text = ocr_service.ocr_image(file_path)
            text = text.strip()

            return ExtractionResult(
                full_text=text,
                page_count=1,
                pages=[PageContent(page_number=1, text=text, ocr_applied=True)],
                ocr_applied=True,
                metadata=metadata
            )
        except OCRNotAvailableError:
            raise
        except Exception as e:
            return ExtractionResult(
                full_text="",
                page_count=1,
                pages=[PageContent(page_number=1, text="", ocr_applied=False)],
                ocr_applied=False,
                metadata=metadata,
                error=f"OCR failed: {str(e)}"
            )
